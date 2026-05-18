from __future__ import annotations

from pathlib import Path

import yaml


ROOT = Path(__file__).resolve().parents[1]
POWERPOINT = ROOT / "powerpoint"
SLIDE_MANIFEST = ROOT / "deck" / "slide_manifest.yaml"
NOTE = POWERPOINT / "README_POWERPOINT_FINALISATION.md"


def append_note(text: str) -> None:
    existing = NOTE.read_text(encoding="utf-8") if NOTE.exists() else "# PowerPoint Finalisation Instructions\n\n"
    if text not in existing:
        NOTE.write_text(existing.rstrip() + "\n\n" + text.strip() + "\n", encoding="utf-8")


def main() -> None:
    POWERPOINT.mkdir(parents=True, exist_ok=True)
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        append_note(
            "## Optional skeleton status\n\n"
            "`python-pptx` is not installed in this environment, so `skeleton_static.pptx` was not generated. "
            "This does not block the asset pipeline; use the storyboard, manifest, slide plates, equations, and MP4s in PowerPoint/Copilot."
        )
        print("python-pptx unavailable; wrote finalisation note")
        return

    data = yaml.safe_load(SLIDE_MANIFEST.read_text(encoding="utf-8")) or {}
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    for slide_spec in data.get("slides", []):
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(14.8), Inches(0.6))
        title_box.text_frame.text = slide_spec["title"]
        image_assets = slide_spec.get("image_assets", [])
        if image_assets:
            image_path = ROOT / image_assets[0]
            if image_path.exists():
                slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.1), width=Inches(14.4))
        note = slide.notes_slide.notes_text_frame
        note.text = slide_spec.get("speaker_note", "")

    output = POWERPOINT / "skeleton_static.pptx"
    prs.save(output)
    print(f"wrote {output}")


if __name__ == "__main__":
    main()
